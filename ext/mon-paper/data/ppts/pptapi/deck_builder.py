from enum import IntEnum
from pathlib import Path
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt

from .animations import add_appear_animation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from lxml import etree

from .types import (
    Content, ContentSlide, EvalSlide, FigureObject, MarkdownList, Object,
    Palette, parse_md_bullets, Slide, Text, TextObject, ThemeMapping, TitleSlide,
)


class Layout(IntEnum):
    TITLE_SLIDE = 0
    TITLE_CONTENT = 1
    COMPARISON = 2
    CONTENT_CAPTION = 3


class DeckBuilder:
    @classmethod
    def _ensure_png(cls, path: Path) -> Path:
        """Convert PDF to PNG if needed."""

        if path.suffix.lower() == ".png":
            # No conversion needed
            return path

        png_path = path.with_suffix(".png")
        do_convert = (not png_path.exists()) or (
            png_path.stat().st_mtime < path.stat().st_mtime
        )
        if do_convert:
            cmd = ["magick", "-density", "300", str(path), str(png_path)]
            subprocess.run(cmd, check=True)

        return png_path

    @classmethod
    def _render_text(cls, tf, content: Content):
        """Render content to a text frame."""

        tf.clear()
        fontsz = content.fontsz

        if isinstance(content, Text):
            p = tf.paragraphs[0]
            p.text = content.text
            if fontsz:
                p.font.size = Pt(fontsz)
        elif isinstance(content, MarkdownList):
            bullets = parse_md_bullets(content.text)
            for i, (level, text) in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = level
                if fontsz:
                    p.font.size = Pt(fontsz)

    def _render_objects(self, slide, objects: list[Object]):
        for obj in objects:
            if isinstance(obj, TextObject):
                self._render_text_object(slide, obj)
            elif isinstance(obj, FigureObject):
                self._render_figure_object(slide, obj)

    def _render_text_object(self, slide, obj: TextObject):
        txbox = slide.shapes.add_textbox(
            Inches(obj.left), Inches(obj.top),
            Inches(obj.width), Inches(obj.height)
        )
        self._render_text(txbox.text_frame, obj.content)

    def _render_figure_object(self, slide, obj: FigureObject):
        img_path = self._ensure_png(obj.path)
        kwargs = {}
        if obj.width:
            kwargs['width'] = Inches(obj.width)
        if obj.height:
            kwargs['height'] = Inches(obj.height)
        slide.shapes.add_picture(str(img_path), Inches(obj.left), Inches(obj.top), **kwargs)

    def __init__(self, template: Path):
        self.prs = Presentation(template)

    def apply_theme(self, palette: Palette, theme: ThemeMapping):
        """Apply palette colors to pptx theme slots."""

        theme_part = self.prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme_xml = etree.fromstring(theme_part.blob)
        clr_scheme = theme_xml.find(".//" + qn("a:clrScheme"))

        for slot in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                     "accent3", "accent4", "accent5", "accent6", "hlink", "folHlink"):
            color_name = getattr(theme, slot)
            hex_color = palette[color_name].lstrip("#")

            color_elem = clr_scheme.find(qn(f"a:{slot}"))
            if color_elem is not None:
                for child in list(color_elem):
                    color_elem.remove(child)
                etree.SubElement(color_elem, qn("a:srgbClr"), val=hex_color)

        theme_part._blob = etree.tostring(theme_xml, xml_declaration=True, encoding="UTF-8", standalone=True)

    def _render_title_slide(self, data: TitleSlide):
        """Title slide: title and subtitle."""

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[Layout.TITLE_SLIDE])
        slide.shapes.title.text = data.title
        if data.subtitle:
            slide.placeholders[1].text = data.subtitle
        self._render_objects(slide, data.objects)

    def _render_content_slide(self, data: ContentSlide):
        """Content slide: title and content"""

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[Layout.TITLE_CONTENT])
        slide.shapes.title.text = data.title
        self._render_text(slide.placeholders[1].text_frame, data.content)
        self._render_objects(slide, data.objects)

    def _render_eval_slide(self, data: EvalSlide):
        """Eval slide: buildout of staged_buildout plots with animations"""

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[Layout.TITLE_CONTENT])
        slide.shapes.title.text = data.title

        # XXX: Do not delete!
        # Slide math: 16 x 9
        # bottomskip: 1.0, topskip: 1.3 => contenth: 6.7
        # leftskip: 0.8, rightskip: 0.8 => contentw: 14.4
        # contentw/2 = 7.2
        # fig: left: 8.2, top: 1.3
        # content: left: 0.8, top: 1.3

        # fig should be 6.5in (w) x 7in (h)

        # Resize content placeholder to left half
        content_ph = slide.placeholders[1]

        # Content dims in inches: (left, top, width, height)
        content_dims = [0.5, 1.4, 7, 6.5]
        cdl, cdt, cdw, cdh = content_dims
        content_ph.left = Inches(cdl)
        content_ph.top = Inches(cdt)
        content_ph.width = Inches(cdw)
        content_ph.height = Inches(cdh)
        self._render_text(content_ph.text_frame, data.bullets)

        # Right side: figures (overlaid, animated)
        # Fig dims: (left, top, width) -- height is scaled, ~7inch
        fig_dims = [8.2, 1.3, 6.5]
        fdl, fdt, fdw = fig_dims
        fdl, fdt, fdw = Inches(fdl), Inches(fdt), Inches(fdw)

        shape_ids = []
        for fig_path in data.figures:
            # python-pptx does not support PDFs. this converts to PNG
            # XXX: this caches, updates may not propagate
            img_path = self._ensure_png(fig_path)
            pic = slide.shapes.add_picture(str(img_path), fdl, fdt, width=fdw)
            shape_ids.append(pic.shape_id)

        if shape_ids:
            add_appear_animation(slide, shape_ids)
        self._render_objects(slide, data.objects)

    def add_slides(self, slides: list[Slide]):
        """Add slide AST to deck."""

        for slide_data in slides:
            if isinstance(slide_data, TitleSlide):
                self._render_title_slide(slide_data)
            elif isinstance(slide_data, ContentSlide):
                self._render_content_slide(slide_data)
            elif isinstance(slide_data, EvalSlide):
                self._render_eval_slide(slide_data)

    def save(self, output: Path):
        """Save deck to file."""

        self.prs.save(output)
        print(f"Saved {output}")
